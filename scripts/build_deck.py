#!/usr/bin/env python3
"""Build docs/runway-software-walkthrough.pptx from the app screenshots.

Crops macOS chrome (menu bar, browser tabs, dock) off the raw screenshots,
then lays out a dark-themed 16:9 deck that explains every module of Runway
using worker W-0046's real computed numbers.

Usage: python3 scripts/build_deck.py
"""

import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(ROOT, "docs", "screenshots")
CROPPED = os.path.join(SHOTS, "cropped")
OUT = os.path.join(ROOT, "docs", "runway-software-walkthrough.pptx")

# ---------- palette (matches the app) ----------
BG = RGBColor(0x0A, 0x0E, 0x13)
PANEL = RGBColor(0x14, 0x1C, 0x25)
BORDER = RGBColor(0x1E, 0x28, 0x33)
TEXT = RGBColor(0xE8, 0xEE, 0xF4)
MUTED = RGBColor(0x8F, 0xA0, 0xAF)
FAINT = RGBColor(0x5C, 0x6B, 0x78)
ACCENT = RGBColor(0xA3, 0xE6, 0x35)
GOOD = RGBColor(0x34, 0xD3, 0x99)
WARN = RGBColor(0xFB, 0xBF, 0x24)
BAD = RGBColor(0xF8, 0x71, 0x71)
INFO = RGBColor(0x60, 0xA5, 0xFA)

FONT = "Helvetica Neue"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Crop fractions measured against the raw 3456x2234 captures:
# macOS menu bar + Chrome tab/URL chrome on top, dock on the bottom.
CROP_TOP = 0.108
CROP_BOTTOM = 0.932


# Shots taken after the first capture carry the macOS floating screenshot
# thumbnail in the bottom-right corner; smear-fill it from the pixels to its left.
THUMBNAIL_SHOTS = {"shot-%02d" % i for i in range(3, 13)}


def erase_thumbnail(img):
    w, h = img.size
    x0, y0 = int(w * 0.872), int(h * 0.825)
    px = img.load()
    for y in range(y0, h):
        fill = px[x0 - 10, y]
        for x in range(x0, w):
            px[x, y] = fill


def crop_shots():
    os.makedirs(CROPPED, exist_ok=True)
    ratios = {}
    for name in sorted(os.listdir(SHOTS)):
        if not name.endswith(".png"):
            continue
        stem = name.replace(".png", "")
        src = os.path.join(SHOTS, name)
        dst = os.path.join(CROPPED, stem + ".jpg")
        img = Image.open(src)
        w, h = img.size
        box = (0, int(h * CROP_TOP), w, int(h * CROP_BOTTOM))
        img = img.crop(box).convert("RGB")
        if stem in THUMBNAIL_SHOTS:
            erase_thumbnail(img)
        if img.width > 2000:
            nh = int(img.height * 2000 / img.width)
            img = img.resize((2000, nh), Image.LANCZOS)
        img.save(dst, "JPEG", quality=88)
        ratios[stem] = img.width / img.height
    return ratios


# ---------- drawing helpers ----------


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def run(par, text, size, color, bold=False, italic=False):
    r = par.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = FONT
    return r


def para(tf, first=False):
    return tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()


def kicker_title(slide, kicker, title, subtitle=None, x=Inches(0.45), w=Inches(12.4)):
    _, tf = textbox(slide, x, Inches(0.32), w, Inches(1.15))
    p = para(tf, first=True)
    run(p, kicker.upper(), 10.5, ACCENT, bold=True)
    p2 = para(tf)
    run(p2, title, 24, TEXT, bold=True)
    if subtitle:
        p3 = para(tf)
        p3.space_before = Pt(2)
        run(p3, subtitle, 11.5, MUTED)


def footer(slide, idx, total):
    _, tf = textbox(slide, Inches(0.45), Inches(7.12), Inches(10.5), Inches(0.3))
    p = para(tf, first=True)
    run(p, "Runway — money copilot for daily earners   ·   kelvine95.github.io/runway-money", 8.5, FAINT)
    _, tf2 = textbox(slide, Inches(12.3), Inches(7.12), Inches(0.6), Inches(0.3))
    p2 = para(tf2, first=True)
    p2.alignment = PP_ALIGN.RIGHT
    run(p2, "%d / %d" % (idx, total), 8.5, FAINT)


def panel(slide, x, y, w, h, fill=PANEL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.045
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def bullets(tf, items, size=11, gap=6):
    for i, (head, body) in enumerate(items):
        p = para(tf, first=(i == 0))
        p.space_after = Pt(gap)
        run(p, "▪  ", size, ACCENT, bold=True)
        if head:
            run(p, head + " — ", size, TEXT, bold=True)
        run(p, body, size, MUTED)


def stat_chip(slide, x, y, w, value, label, color=TEXT):
    panel(slide, x, y, w, Inches(0.95))
    _, tf = textbox(slide, x + Inches(0.16), y + Inches(0.12), w - Inches(0.32), Inches(0.75))
    p = para(tf, first=True)
    run(p, value, 16, color, bold=True)
    p2 = para(tf)
    run(p2, label, 8.5, MUTED)


# ---------- slide builders ----------


def screenshot_slide(prs, ratios, idx, total, shot, kicker, title, subtitle, seeing, hood, hood_title="Under the hood"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    kicker_title(slide, kicker, title, subtitle)

    img_path = os.path.join(CROPPED, shot + ".jpg")
    ratio = ratios[shot]
    img_w = Inches(8.55)
    img_h = Emu(int(img_w / ratio))
    max_h = Inches(5.35)
    if img_h > max_h:
        img_h = max_h
        img_w = Emu(int(img_h * ratio))
    pic = slide.shapes.add_picture(img_path, Inches(0.45), Inches(1.55), width=img_w, height=img_h)
    pic.line.color.rgb = BORDER
    pic.line.width = Pt(1)

    px = Inches(9.3)
    pw = Inches(3.6)
    _, tf = textbox(slide, px, Inches(1.55), pw, Inches(0.3))
    p = para(tf, first=True)
    run(p, "WHAT YOU'RE SEEING", 9.5, ACCENT, bold=True)
    _, tf2 = textbox(slide, px, Inches(1.92), pw, Inches(3.1))
    bullets(tf2, seeing, size=10.5, gap=5)

    if hood:
        hy = Inches(5.15)
        hh = Inches(1.78)
        panel(slide, px - Inches(0.14), hy, pw + Inches(0.28), hh)
        _, tf3 = textbox(slide, px, hy + Inches(0.12), pw, hh - Inches(0.24))
        p = para(tf3, first=True)
        run(p, hood_title.upper(), 9, INFO, bold=True)
        p2 = para(tf3)
        p2.space_before = Pt(3)
        run(p2, hood, 9.5, MUTED)

    footer(slide, idx, total)


def build():
    ratios = crop_shots()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 17

    # ---- 1 · Title ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    mark = panel(s, Inches(0.75), Inches(1.5), Inches(0.85), Inches(0.85), fill=ACCENT)
    _, tfm = textbox(s, Inches(0.75), Inches(1.62), Inches(0.85), Inches(0.6))
    pm = para(tfm, first=True)
    pm.alignment = PP_ALIGN.CENTER
    run(pm, "R", 30, BG, bold=True)
    _, tf = textbox(s, Inches(1.85), Inches(1.42), Inches(10.5), Inches(1.1))
    p = para(tf, first=True)
    run(p, "Runway", 44, TEXT, bold=True)
    p2 = para(tf)
    run(p2, "A money copilot for people who get paid every day — not twice a month", 15, MUTED)

    _, tf3 = textbox(s, Inches(0.78), Inches(3.1), Inches(11.5), Inches(0.9))
    p = para(tf3, first=True)
    run(p, "SOFTWARE WALKTHROUGH · HACKATHON BUILD · JULY 2026", 11, ACCENT, bold=True)
    p2 = para(tf3)
    p2.space_before = Pt(6)
    run(p2, "Live demo:  ", 12, MUTED)
    run(p2, "kelvine95.github.io/runway-money", 12, TEXT, bold=True)
    run(p2, "      Code:  ", 12, MUTED)
    run(p2, "github.com/kelvine95/runway-money", 12, TEXT, bold=True)

    chips = [
        ("6 modules", "Today · Life Pulse · Steady Paycheck · Bill Runway · Shift ROI · Advance Audit", TEXT),
        ("100% client-side", "React + TypeScript, no backend, free GitHub Pages hosting", GOOD),
        ("220 workers", "synthetic Alberta cohort: 12,204 shifts · 31,726 transactions", INFO),
        ("535 advances", "$1,435 in wage-advance fees observed in ~14 weeks", BAD),
    ]
    cx = Inches(0.78)
    for value, label, color in chips:
        stat_chip(s, cx, Inches(4.55), Inches(2.95), value, label, color)
        cx += Inches(3.08)
    footer(s, 1, total)

    # ---- 2 · What the software is ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    kicker_title(
        s,
        "Overview",
        "What Runway does",
        "It answers the question a bank app can't: “Given how I actually earn, can I spend today, will rent clear, and what should I do about it?”",
    )
    _, tf = textbox(s, Inches(0.45), Inches(1.75), Inches(5.4), Inches(4.9))
    bullets(
        tf,
        [
            ("Ingest", "six CSV feeds per worker — profile, daily shift earnings, bank transactions, recurring bills, wage advances, weekly summaries."),
            ("Rebuild", "a clean balance ledger from transaction deltas (the raw feed's running balances are inconsistent, so the app never trusts them blindly)."),
            ("Model", "each worker: income distribution, income floor (their own 25th-percentile week), essential spend rate, bill calendar, median take-home per shift."),
            ("Decide", "six views that speak in the worker's units — days of runway, shifts needed, funded confidence — instead of abstract budget categories."),
            ("Respect", "an ethical engagement loop: nudges only on surplus days, never savings pressure on lean days."),
        ],
        size=11.5,
        gap=9,
    )
    modules = [
        ("1 · Today", "Safe-to-spend after bills, essentials and a 7-day buffer — under a bad-week income assumption.", TEXT),
        ("2 · Steady Paycheck", "Pay-yourself-a-salary automation, replayed and proven against real history.", GOOD),
        ("3 · Bill Runway", "400-trial Monte Carlo of the next 30 days; per-bill funding confidence in shifts.", INFO),
        ("4 · Shift ROI", "Take-home $/hour by weekday, shift type and employer — the income lever.", TEXT),
        ("5 · Advance Audit", "Every wage advance priced as an effective APR, timed against rent day.", BAD),
        ("6 · Life Pulse", "Personal inflation, living-wage gap, CA/US baselines + daily check-in.", WARN),
    ]
    my = Inches(1.75)
    for name, desc, color in modules:
        panel(s, Inches(6.15), my, Inches(6.7), Inches(0.72))
        _, tfm = textbox(s, Inches(6.35), my + Inches(0.09), Inches(6.3), Inches(0.56))
        p = para(tfm, first=True)
        run(p, name + "   ", 10.5, color, bold=True)
        run(p, desc, 9.5, MUTED)
        my += Inches(0.82)
    footer(s, 2, total)

    # ---- 3 · Demo worker ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    kicker_title(
        s,
        "The demo persona",
        "Meet W-0046 — landscaping & grounds, Calgary",
        "Every number in the next twelve slides is computed live in the browser from this worker's raw records. Nothing is mocked.",
    )
    facts = [
        ("Daily pay", "45% income volatility", WARN),
        ("Severe rent burden", "1 dependent at home", BAD),
        ("54 shifts on record", "~3.9 workdays per week", TEXT),
        ("$206 median shift", "≈ 8.7 hours · $24.57/hr", GOOD),
        ("8 wage advances", "$40.44 in fees in 96 days", BAD),
        ("$2,245 in the bank", "…and still $0 safe to spend", INFO),
    ]
    fx, fy = Inches(0.6), Inches(2.0)
    for i, (value, label, color) in enumerate(facts):
        stat_chip(s, fx + Inches(4.15) * (i % 3), fy + Inches(1.25) * (i // 3), Inches(3.9), value, label, color)
    panel(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.7))
    _, tf = textbox(s, Inches(0.85), Inches(5.0), Inches(11.6), Inches(1.35))
    p = para(tf, first=True)
    run(p, "WHY THIS PERSONA MATTERS   ", 10, ACCENT, bold=True)
    p2 = para(tf)
    p2.space_before = Pt(4)
    run(
        p2,
        "W-0046 looks fine in a banking app: $2,245 in the account. But childcare ($805) hits tomorrow and rent ($1,956) in nine days, "
        "while income arrives in unpredictable $100–$320 daily chunks. This is exactly the gap between a balance and a plan — "
        "the worker has already paid $40.44 in advance fees this quarter renting liquidity they could have owned.",
        11,
        MUTED,
    )
    footer(s, 3, total)

    # ---- 4-15 · module screenshot slides ----
    screenshot_slide(
        prs, ratios, 4, total, "shot-01",
        "Module 1 · Today", "Safe to Spend — the headline number",
        "The opening screen. Not the bank balance — what is genuinely spendable today.",
        [
            ("$0.00 safe today", "even though the bank shows $2,245. The gap is the whole product."),
            ("Four cards", "safe-to-spend · raw bank balance · 20 days of zero-income runway · next bill (childcare, $805, due tomorrow)."),
            ("Red alert", "worker is $302 short of covering 14 days + a 7-day buffer ≈ 1.5 extra shifts at their $206 median."),
            ("Balance chart", "8 weeks of the rebuilt ledger — the sawtooth cliffs are bill days."),
        ],
        "Flex = balance + conservative income − bills due − essentials − buffer floor. "
        "Safe today = max(0, Flex) ÷ 14. Conservative income uses the worker's own 25th-percentile week — planning for a bad stretch, not an average one.",
    )

    screenshot_slide(
        prs, ratios, 5, total, "shot-02",
        "Module 1 · Today", "How the number is built, line by line",
        "Full transparency: the same screen shows the math and the 35-day bill outlook.",
        [
            ("The breakdown", "+$2,245 balance, +$1,245 conservative income, −$2,808 bills, −$195.96 essentials ($14.00/day observed), −$787.01 buffer floor → −$301.72 flex."),
            ("Every bill scored", "childcare $805 → 100% funded confidence · rent $1,956 → 97% · phone $47 → 98% · utilities beyond forecast."),
            ("No black box", "each reservation line says what it is and why it is held back."),
        ],
        "Funding confidence comes from 400 simulated futures (see Bill Runway). Essentials are the trailing-56-day variable essential spend; "
        "the buffer floor is 7 days of essentials + prorated monthly obligations.",
    )

    screenshot_slide(
        prs, ratios, 6, total, "shot-05",
        "Module 2 · Steady Paycheck", "Turning volatile pay into a salary",
        "The classic irregular-income method — computed and proven instead of preached.",
        [
            ("$89/day paycheck", "$622/week — the 25th percentile of this worker's own observed full weeks."),
            ("98% reliability", "replaying 96 real days, the holding buffer could pay the full $89 on 98% of days."),
            ("$2,808 buffer built", "surplus from strong days sits between the worker and the next slow week."),
            ("$40.44 fees avoided", "8 of 8 wage advances would never have been needed."),
        ],
        "Replay: every real earning lands in a holding buffer; the worker is paid a fixed daily amount from it. "
        "An advance is marked avoidable when the buffer already held enough on the day it was requested.",
    )

    screenshot_slide(
        prs, ratios, 7, total, "shot-06",
        "Module 2 · Steady Paycheck", "The buffer that makes it possible",
        "Hovering the chart shows the mechanism on a real zero-income day.",
        [
            ("Tue, Apr 28", "actual earnings $0.00 — steady paycheck still pays $89.00. That is the smoothing in action."),
            ("Buffer curve", "grows from $0 to ~$2,800 over three months; slow days drain it, strong days refill it."),
            ("The insight", "this worker never needed to borrow — they needed their own strong days redistributed by a few weeks."),
        ],
        "Floor = P25(weekly net) ÷ 7, so roughly 3 of every 4 weeks build surplus by construction. "
        "The buffer is liquidity the worker owns instead of renting from an advance app.",
    )

    screenshot_slide(
        prs, ratios, 8, total, "shot-07",
        "Module 3 · Bill Runway", "400 possible versions of the next 30 days",
        "Instead of one guess, a distribution — sampled from how this worker actually earns.",
        [
            ("Median path: never breaks", "the typical simulated future stays above $0 for the full 30 days."),
            ("0% tail risk", "share of simulations ending the month negative."),
            ("Fan chart", "shaded band = 10th–90th percentile of 400 simulations; line = median; dots = bill due dates."),
            ("$2,808 of bills", "3 scheduled payments inside the window."),
        ],
        "Each trial samples daily income from the trailing 56-day pool (off-days included as $0), subtracts daily essentials, "
        "and pays bills on their due dates. Seeded per worker, so results are stable across reloads.",
    )

    screenshot_slide(
        prs, ratios, 9, total, "shot-08",
        "Module 3 · Bill Runway", "Per-bill funding confidence, in shifts",
        "Every due date gets its own probability — and gaps are quoted in the worker's unit.",
        [
            ("Hover any day", "Jul 13: 10th–90th percentile $1,952–$2,728, median $2,365."),
            ("Childcare 100%", "$805 due in 1 day — fully funded in every simulated future."),
            ("Rent 97%", "$1,956 in 9 days — “needs ~0.3 extra shifts if the gap hits.”"),
            ("Why shifts?", "dollars describe the problem; shifts describe the action. Picking up work is the lever this user controls."),
        ],
        "Funding confidence = share of trials where the bill clears without the balance going negative. "
        "Shifts needed = median shortfall across failing trials ÷ median net per shift.",
    )

    screenshot_slide(
        prs, ratios, 10, total, "shot-09",
        "Module 4 · Shift ROI", "Which work actually pays",
        "Budget apps optimize spending. For a daily earner the bigger lever is which shifts to say yes to.",
        [
            ("$24.57/hr median", "take-home across all 54 recorded shifts."),
            ("Wednesday wins", "$27.07/hr median vs Thursday's $21.39/hr — highlighted lime vs red in the chart."),
            ("$198/month insight", "swapping one Thu shift a week for a Wed shift ≈ 1.0 extra shift of money, without working it."),
            ("39% same-day pay", "how often this worker is paid the day they work — a liquidity fact, not just a wage fact."),
        ],
        "Median net ÷ hours per slot, split by weekday and shift type. Only slots with ≥5 recorded shifts are treated as credible enough to recommend.",
    )

    screenshot_slide(
        prs, ratios, 11, total, "shot-10",
        "Module 4 · Shift ROI", "By employer — who respects your hour",
        "Multi-employer income is normal for daily earners. Runway ranks the payers.",
        [
            ("EMP-016 is the anchor", "43 of 54 shifts, $24.88/hr median, $205.60 per shift, 40% paid same-day."),
            ("One-off gigs ranked", "from EMP-033 at $28.49/hr down to EMP-038 at $15.62/hr — a $13/hr spread for the same worker."),
            ("Columns that matter", "shifts (sample size), $/hr, $/shift, tips share, same-day-pay share."),
        ],
        "This table is the seed of a marketplace insight: the same person, same city, same quarter — "
        "and nearly 2× difference in effective wage depending on who they work for.",
    )

    screenshot_slide(
        prs, ratios, 12, total, "shot-11",
        "Module 5 · Advance Audit", "Pricing the cash-advance habit honestly",
        "What borrowing against your own wages actually costs, in numbers lenders don't show.",
        [
            ("8 advances / $721.42", "borrowed over 96 days, costing $40.44 in fees ≈ $154/year at this pace."),
            ("261% median effective APR", "fee ÷ amount, annualized over days outstanding."),
            ("38% taken in rent week", "the timing histogram shows borrowing clustering at the bill cliff."),
            ("The verdict", "8 of 8 advances happened on days the steady-paycheck buffer would have covered — that $40.44 bought liquidity this worker could have owned free."),
        ],
        "Effective APR = (fee ÷ amount) × (365 ÷ days outstanding). Reasons are taken from the advance request records: "
        "emergency, groceries, rent gap, bill due, childcare — timing failures, not luxury overspend.",
    )

    screenshot_slide(
        prs, ratios, 13, total, "shot-12",
        "Module 5 · Advance Audit", "Every advance itemized — and the cohort view",
        "A per-advance ledger for the worker, then the zoom-out across all 220 workers.",
        [
            ("The ledger", "each advance with date, reason, amount, fee, days outstanding, effective APR, repayment status."),
            ("Small fees, huge rates", "a $48.60 childcare advance with a $1.99 fee repaid next day = 1,495% effective APR."),
            ("Cohort strip", "128 of 220 workers advanced · 535 advances · $1,435 fees ≈ $5,455/year pace · $11.21 per advance-user."),
            ("Why it matters", "a quiet tax concentrated on the people with the least slack — and a measurable target for the product to eliminate."),
        ],
        None,
    )

    screenshot_slide(
        prs, ratios, 14, total, "shot-03",
        "Module 6 · Life Pulse", "A daily ritual that respects lean days",
        "The retention surface: check in, see your place in the economy, take one small win.",
        [
            ("Energy check-in", "“How much energy do you have — not how productive should you be?” Low battery / Steady / Ready to go."),
            ("Personal inflation 2.2%", "official category inflation weighted by where W-0046 actually spends."),
            ("Living-wage gap", "$23.75/hr observed take-home vs Calgary's $26.50/hr benchmark → the callout names it a structural gap, not a budgeting failure."),
            ("Adaptive missions", "because there is no safe flex today, the app refuses to manufacture a savings goal — today's wins are protection and rest."),
        ],
        "Resilience points and a gentle streak live in localStorage. No lost-streak shame, no borrow nudges — "
        "surplus days invite a tiny buffer move; lean days never apply savings pressure.",
    )

    screenshot_slide(
        prs, ratios, 15, total, "shot-04",
        "Module 6 · Life Pulse", "The economy layer — context, not hindsight",
        "Official macro baselines personalized to the worker, with sources linked in-app.",
        [
            ("Weighted inflation chart", "transit inflation (+6.7%) dominates for this landscaper because that is where their essential dollars go."),
            ("Canada vs US pulse", "CPI 2.8% vs 3.5% · Alberta CPI 3.4% · Alberta unemployment 7.0% · BoC 2.25% vs Fed 3.50–3.75%."),
            ("The baseline strip", "$15.00/hr Alberta minimum wage unchanged since 2018 · groceries +3.9% · gasoline +20.5%."),
            ("Design principle", "macro data is displayed as a dated snapshot; it is never injected into the historical forecast, avoiding fake precision."),
        ],
        None,
    )

    # ---- 16 · Under the hood ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    kicker_title(
        s,
        "Engineering",
        "Under the hood",
        "Small, legible, and deliberately honest about its assumptions.",
    )
    cols = [
        ("STACK", [
            ("Frontend", "Vite + React + TypeScript, Recharts for charts, PapaParse for CSV ingestion."),
            ("Hosting", "static build on GitHub Pages — zero backend, zero hosting cost."),
            ("Analytics", "every model runs client-side in TypeScript: safeToSpend, steadyPaycheck, runway (Monte Carlo), shiftRoi, advances, economicContext."),
            ("Determinism", "Monte Carlo is seeded per worker (mulberry32), so numbers are stable across reloads and demos."),
        ]),
        ("DATA & METHOD", [
            ("Six feeds", "workers (220) · daily earnings (12,204) · transactions (31,726) · obligations (849) · advances (535) · weekly summaries."),
            ("Ledger rebuild", "canonical balances recomputed from transaction deltas — the raw running-balance field is inconsistent and never trusted."),
            ("Conservatism", "income floors use P25, not averages; averages overstate capacity for volatile earners and recreate the cliff."),
            ("Stated limits", "synthetic cohort; i.i.d. daily sampling (no weekday seasonality yet); P25 floor and 7-day buffer are defaults, not yet personalized."),
        ]),
    ]
    cx = Inches(0.45)
    for heading, items in cols:
        _, tfh = textbox(s, cx, Inches(1.8), Inches(6.1), Inches(0.3))
        p = para(tfh, first=True)
        run(p, heading, 10.5, INFO, bold=True)
        _, tfc = textbox(s, cx, Inches(2.2), Inches(6.1), Inches(4.4))
        bullets(tfc, items, size=10.5, gap=8)
        cx += Inches(6.4)
    footer(s, 16, total)

    # ---- 17 · Try it ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    kicker_title(
        s,
        "Try it now",
        "Live demo, open code",
        None,
    )
    rows = [
        ("Live demo", "https://kelvine95.github.io/runway-money/", "Opens on the cohort's heaviest advance user; switch among all 220 workers in the sidebar."),
        ("Source code", "https://github.com/kelvine95/runway-money", "MIT-style hackathon repo — data, analytics engine, UI, and this deck."),
        ("Run locally", "npm install && npm run dev", "No API keys, no configuration. The dataset ships with the repo."),
    ]
    ry = Inches(1.9)
    for label, big, small in rows:
        panel(s, Inches(0.6), ry, Inches(12.1), Inches(1.15))
        _, tfr = textbox(s, Inches(0.9), ry + Inches(0.14), Inches(11.6), Inches(0.9))
        p = para(tfr, first=True)
        run(p, label.upper() + "   ", 10, ACCENT, bold=True)
        run(p, big, 14, TEXT, bold=True)
        p2 = para(tfr)
        p2.space_before = Pt(3)
        run(p2, small, 10, MUTED)
        ry += Inches(1.35)
    panel(s, Inches(0.6), ry, Inches(12.1), Inches(0.9), fill=BG)
    _, tfq = textbox(s, Inches(0.9), ry + Inches(0.16), Inches(11.6), Inches(0.6))
    p = para(tfq, first=True)
    run(
        p,
        "“The market already charges workers for surviving income volatility. Runway's bet is that preventing the cliff "
        "is a better product than renting people their own wages.”",
        12.5,
        TEXT,
        italic=True,
    )
    footer(s, 17, total)

    prs.save(OUT)
    print("Saved %s (%.1f MB)" % (OUT, os.path.getsize(OUT) / 1e6))


if __name__ == "__main__":
    build()
